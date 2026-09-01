"""
Spine Engine — PowerPoint Parser.

Parses PowerPoint presentations (.pptx) — typically training decks.
Uses python-pptx for real parsing.
"""

from __future__ import annotations

import io
import logging

from nexus_sdk.events import fire_stub_alert

from ..models import ExtractedTable

logger = logging.getLogger("nexus.spine.parsers.powerpoint")

# Module-level event bus and stub tracking
_event_bus = None
_stub_counts: dict[str, int] = {}


def set_event_bus(bus) -> None:
    """Called by SpineEngine.on_startup() to inject the event bus reference."""
    global _event_bus
    _event_bus = bus


class PowerPointParser:
    """Parse PowerPoint presentations (training decks)."""

    @staticmethod
    def parse(content: bytes, document_id: str) -> dict:
        try:
            from pptx import Presentation  # type: ignore[import-not-found]
            return PowerPointParser._parse_with_pptx(content, document_id)
        except ImportError:
            return PowerPointParser._stub_parse(content, document_id)

    @staticmethod
    def _parse_with_pptx(content: bytes, document_id: str) -> dict:
        from pptx import Presentation  # type: ignore[import-not-found]

        prs = Presentation(io.BytesIO(content))
        slides = []
        full_text = ""
        tables = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = ""
            title = ""

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text += text + "\n"
                            if not title and shape.shape_id == slide.shapes[0].shape_id:
                                title = text

                # Extract tables from slides
                if shape.has_table:
                    tbl = shape.table
                    rows = []
                    for row in tbl.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        headers = rows[0]
                        tables.append(ExtractedTable(
                            document_id=document_id,
                            page_number=slide_num,
                            headers=headers,
                            rows=rows[1:],
                            row_count=len(rows) - 1,
                            col_count=len(headers),
                            metadata={"slide_number": slide_num},
                        ))

            slides.append({
                "slide_number": slide_num,
                "title": title,
                "text": slide_text,
                "char_count": len(slide_text),
            })
            full_text += f"--- Slide {slide_num}: {title} ---\n{slide_text}\n"

        return {
            "slides": slides,
            "tables": tables,
            "full_text": full_text,
            "page_count": len(slides),
        }

    @staticmethod
    def _stub_parse(content: bytes, document_id: str) -> dict:
        _stub_counts["pptx"] = _stub_counts.get("pptx", 0) + 1
        logger.warning("spine: PowerPoint parser stub fallback #%d", _stub_counts["pptx"])
        fire_stub_alert(
            _event_bus, "spine", "pptx_parser",
            fallback_count=_stub_counts["pptx"],
            reason="python-pptx not installed",
        )
        return {
            "slides": [],
            "tables": [],
            "full_text": "[Stub PowerPoint — install python-pptx for real parsing]",
            "page_count": 1,
        }
